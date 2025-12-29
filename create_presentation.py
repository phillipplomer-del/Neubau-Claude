#!/usr/bin/env python3
"""
Galadriel Project Presentation Generator
Erstellt eine professionelle PowerPoint-Präsentation für das Galadriel Business Intelligence Dashboard
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# Farben definieren
BLUE_PRIMARY = RGBColor(30, 58, 138)      # Dunkelblau
BLUE_SECONDARY = RGBColor(59, 130, 246)   # Hellblau
GRAY_DARK = RGBColor(31, 41, 55)          # Dunkelgrau
GRAY_LIGHT = RGBColor(107, 114, 128)      # Hellgrau
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(34, 197, 94)             # Grün für Erfolg
ORANGE = RGBColor(249, 115, 22)           # Orange für Warnung


def add_title_slide(prs):
    """Titelfolie erstellen"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Hintergrund-Shape (Blauer Balken oben)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE_PRIMARY
    shape.line.fill.background()

    # Logo/Icon Bereich
    icon_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6), Inches(0.5), Inches(1.33), Inches(1.33))
    icon_shape.fill.solid()
    icon_shape.fill.fore_color.rgb = WHITE
    icon_shape.line.fill.background()

    # Haupttitel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.33), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "GALADRIEL"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Untertitel
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.33), Inches(1))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Business Intelligence Dashboard"
    p.font.size = Pt(32)
    p.font.color.rgb = BLUE_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    # Beschreibung
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.33), Inches(1.5))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Integrierte Lösung für Sales, Produktion & Projektmanagement"
    p.font.size = Pt(24)
    p.font.color.rgb = GRAY_DARK
    p.alignment = PP_ALIGN.CENTER

    # Datum/Version
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Dezember 2025"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY_LIGHT
    p.alignment = PP_ALIGN.CENTER


def add_agenda_slide(prs):
    """Agenda-Folie erstellen"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    add_header(slide, "Agenda")

    agenda_items = [
        ("01", "Projektübersicht", "Was ist Galadriel?"),
        ("02", "Technologie-Stack", "Moderne Architektur & Tools"),
        ("03", "Module & Features", "Sales, Produktion, Projektmanagement"),
        ("04", "Datenarchitektur", "IndexedDB & Firebase Integration"),
        ("05", "Demo & Screenshots", "Live-Einblicke"),
        ("06", "Zusammenfassung", "Vorteile & nächste Schritte")
    ]

    y_start = 1.8
    for i, (num, title, desc) in enumerate(agenda_items):
        # Nummer
        num_box = slide.shapes.add_textbox(Inches(1), Inches(y_start + i * 0.8), Inches(0.8), Inches(0.6))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = BLUE_SECONDARY

        # Titel
        title_box = slide.shapes.add_textbox(Inches(2), Inches(y_start + i * 0.8), Inches(4), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = GRAY_DARK

        # Beschreibung
        desc_box = slide.shapes.add_textbox(Inches(2), Inches(y_start + 0.35 + i * 0.8), Inches(8), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = GRAY_LIGHT


def add_header(slide, title):
    """Standard-Header für Folien"""
    # Blaue Linie oben
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.1))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE_PRIMARY
    line.line.fill.background()

    # Titel
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = BLUE_PRIMARY


def add_overview_slide(prs):
    """Projektübersicht-Folie"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_header(slide, "Projektübersicht")

    # Haupttext
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(1))
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Galadriel ist ein modernes Business Intelligence Dashboard zur Verwaltung und Überwachung von Geschäftsdaten aus drei Kernbereichen."
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY_DARK

    # Drei Hauptbereiche als Karten
    areas = [
        ("Sales", "Offene Lieferungen\nKundenaufträge\nLieferstatus-Tracking\nKPI-Überwachung", BLUE_SECONDARY),
        ("Produktion", "Produktionsplanung\nSoll-Ist Vergleiche\nZeitplanung\nRessourcenverwaltung", GREEN),
        ("Projektmanagement", "Projekt-Controlling\nBudget-Überwachung\nProfitabilitätsanalyse\nTrend-Analysen", ORANGE)
    ]

    x_positions = [0.5, 4.6, 8.7]
    for i, (title, content, color) in enumerate(areas):
        # Karten-Hintergrund
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(x_positions[i]), Inches(2.8),
                                      Inches(3.8), Inches(3.5))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(248, 250, 252)
        card.line.color.rgb = color
        card.line.width = Pt(2)

        # Farbiger Header-Balken
        header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           Inches(x_positions[i]), Inches(2.8),
                                           Inches(3.8), Inches(0.6))
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = color
        header_bar.line.fill.background()

        # Titel
        title_box = slide.shapes.add_textbox(Inches(x_positions[i] + 0.1), Inches(2.85),
                                            Inches(3.6), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Content
        content_box = slide.shapes.add_textbox(Inches(x_positions[i] + 0.2), Inches(3.6),
                                              Inches(3.4), Inches(2.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        for line in content.split('\n'):
            p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
            p.text = "• " + line
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY_DARK
            p.space_before = Pt(8)


def add_tech_stack_slide(prs):
    """Technologie-Stack-Folie"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_header(slide, "Technologie-Stack")

    # Technologie-Kategorien
    categories = [
        ("Frontend", ["React 18.3 + TypeScript 5.7", "Vite 6.0 (Build Tool)", "TailwindCSS 3.4"], Inches(0.5)),
        ("UI/Charts", ["Recharts (Datenvisualisierung)", "TanStack Table (Tabellen)", "Frappe-Gantt (Zeitplanung)"], Inches(4.5)),
        ("Daten", ["IndexedDB (Offline-First)", "Firebase Firestore (Cloud)", "XLSX Parser (Excel-Import)"], Inches(8.5))
    ]

    for title, items, x_pos in categories:
        # Kategorie-Titel
        cat_box = slide.shapes.add_textbox(x_pos, Inches(1.6), Inches(3.5), Inches(0.5))
        tf = cat_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = BLUE_PRIMARY

        # Items
        for i, item in enumerate(items):
            item_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                               x_pos, Inches(2.2 + i * 0.7),
                                               Inches(3.8), Inches(0.55))
            item_shape.fill.solid()
            item_shape.fill.fore_color.rgb = RGBColor(239, 246, 255)
            item_shape.line.color.rgb = BLUE_SECONDARY

            item_box = slide.shapes.add_textbox(x_pos + Inches(0.15), Inches(2.3 + i * 0.7),
                                               Inches(3.5), Inches(0.4))
            tf = item_box.text_frame
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(14)
            p.font.color.rgb = GRAY_DARK

    # Zusätzliche Tools
    tools_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(0.4))
    tf = tools_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Weitere Technologien:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GRAY_DARK

    additional_tools = [
        "React Hook Form + Zod (Formular-Validierung)",
        "jsPDF + HTML2Canvas (PDF-Export)",
        "Date-fns (Datumsverarbeitung)",
        "Lucide React (Icon-Bibliothek)"
    ]

    for i, tool in enumerate(additional_tools):
        col = i % 2
        row = i // 2
        tool_box = slide.shapes.add_textbox(Inches(0.5 + col * 6), Inches(5.0 + row * 0.5),
                                           Inches(5.5), Inches(0.4))
        tf = tool_box.text_frame
        p = tf.paragraphs[0]
        p.text = "• " + tool
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY_LIGHT


def add_sales_module_slide(prs):
    """Sales-Modul-Folie"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_header(slide, "Sales-Modul")

    # Beschreibung
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(0.6))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Zentrales Dashboard zur Überwachung aller offenen Lieferungen und Kundenaufträge"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY_LIGHT

    # Features in zwei Spalten
    left_features = [
        ("Dashboard & KPIs", "Übersicht mit Gesamtwert, verzögerten Lieferungen und Status-Verteilung"),
        ("Filterbare Tabellen", "Nach Artikelnummer, Projektnummer, Datum und Freitext filterbar"),
        ("Gruppierte Ansicht", "Projekte gruppiert für 'Beobachtete' Lieferungen"),
        ("PDF-Export", "Automatische KPI-Zusammenfassung als PDF")
    ]

    right_features = [
        ("Kommentar-System", "Status-Markierung: Kritisch, Gefährdet, Beobachtet"),
        ("Firebase-Sync", "Echtzeit-Synchronisation aller Kommentare"),
        ("Status-Updates", "Live-Tracking kritischer Lieferungen"),
        ("Sortierung", "Flexible Sortierung nach allen Spalten")
    ]

    # Linke Spalte
    for i, (title, desc) in enumerate(left_features):
        y_pos = 2.0 + i * 1.1

        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(5.8), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "✓ " + title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLUE_PRIMARY

        desc_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.35), Inches(5.6), Inches(0.5))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY_DARK

    # Rechte Spalte
    for i, (title, desc) in enumerate(right_features):
        y_pos = 2.0 + i * 1.1

        title_box = slide.shapes.add_textbox(Inches(6.8), Inches(y_pos), Inches(5.8), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "✓ " + title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLUE_PRIMARY

        desc_box = slide.shapes.add_textbox(Inches(7.0), Inches(y_pos + 0.35), Inches(5.6), Inches(0.5))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY_DARK


def add_controlling_module_slide(prs):
    """Controlling-Modul-Folie"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_header(slide, "Projektmanagement & Controlling")

    # Beschreibung
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(0.6))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Umfassende Analyse-Tools für Projekt-Performance und Finanzübersicht"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY_LIGHT

    # KPI-Cards Darstellung
    kpis = [
        ("Projekte", "Gesamtübersicht\naller Projekte", BLUE_SECONDARY),
        ("Budget", "Budget vs.\nActual Kosten", GREEN),
        ("ROI", "Return on\nInvestment", ORANGE),
        ("Gewinn", "Gewinnmargen-\nBerechnung", RGBColor(168, 85, 247))
    ]

    x_positions = [0.5, 3.4, 6.3, 9.2]
    for i, (title, desc, color) in enumerate(kpis):
        # KPI Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(x_positions[i]), Inches(2.0),
                                     Inches(2.7), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(x_positions[i] + 0.1), Inches(2.1),
                                            Inches(2.5), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(Inches(x_positions[i] + 0.1), Inches(2.6),
                                           Inches(2.5), Inches(0.8))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Chart-Typen
    charts_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12), Inches(0.4))
    tf = charts_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Verfügbare Visualisierungen:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GRAY_DARK

    charts = [
        ("Zeitreihen-Charts", "Trend-Analyse über Monate/Jahre"),
        ("Bar-Charts", "Budget vs. Actual Vergleich"),
        ("Pie-Charts", "Verteilung nach Manager & Kategorie"),
        ("PDF-Reports", "Exportierbare Controlling-Berichte")
    ]

    for i, (title, desc) in enumerate(charts):
        col = i % 2
        row = i // 2

        # Icon-Shape
        icon = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(0.5 + col * 6.3), Inches(4.4 + row * 1.2),
                                     Inches(0.4), Inches(0.4))
        icon.fill.solid()
        icon.fill.fore_color.rgb = BLUE_SECONDARY
        icon.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(1.0 + col * 6.3), Inches(4.35 + row * 1.2),
                                            Inches(5), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLUE_PRIMARY

        # Description
        desc_box = slide.shapes.add_textbox(Inches(1.0 + col * 6.3), Inches(4.7 + row * 1.2),
                                           Inches(5), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY_LIGHT


def add_data_comparison_slide(prs):
    """Datenabgleich-Folie"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_header(slide, "Datenabgleich")

    # Beschreibung
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(0.6))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Cross-Department-Vergleich für Datenqualität und Konsistenzprüfung"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY_LIGHT

    # Venn-Diagramm Darstellung (vereinfacht mit Shapes)
    # Drei überlappende Kreise mit helleren Farben für bessere Lesbarkeit
    circles = [
        (Inches(3.5), Inches(2.5), "Sales", BLUE_SECONDARY, RGBColor(191, 219, 254)),
        (Inches(5.5), Inches(2.5), "Produktion", GREEN, RGBColor(187, 247, 208)),
        (Inches(4.5), Inches(4.0), "Projekt-\nmanagement", ORANGE, RGBColor(254, 215, 170))
    ]

    for x, y, label, border_color, fill_color in circles:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(2.5), Inches(2.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = fill_color
        circle.line.color.rgb = border_color
        circle.line.width = Pt(2)

        label_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.9),
                                            Inches(1.9), Inches(0.8))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = GRAY_DARK
        p.alignment = PP_ALIGN.CENTER

    # Features rechts
    features_title = slide.shapes.add_textbox(Inches(8.5), Inches(2.0), Inches(4), Inches(0.4))
    tf = features_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Funktionen:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GRAY_DARK

    features = [
        "Projektnummern-Vergleich",
        "Artikelnummern-Matching",
        "Überlappungs-Filter",
        "Visuelle Status-Indikatoren",
        "Datenqualitäts-Analyse"
    ]

    for i, feature in enumerate(features):
        feature_box = slide.shapes.add_textbox(Inches(8.5), Inches(2.5 + i * 0.5),
                                              Inches(4), Inches(0.4))
        tf = feature_box.text_frame
        p = tf.paragraphs[0]
        p.text = "✓ " + feature
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY_DARK


def add_architecture_slide(prs):
    """Architektur-Folie"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_header(slide, "Datenarchitektur")

    # Beschreibung
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12), Inches(0.6))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Offline-First Architektur mit Cloud-Synchronisation"
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY_LIGHT

    # Architektur-Schichten
    layers = [
        ("Excel Import", "XLSX-Parser mit Auto-Detection", Inches(1), RGBColor(254, 243, 199)),
        ("Validierung", "Zod Schema-Validierung", Inches(2.2), RGBColor(254, 226, 226)),
        ("IndexedDB", "Lokale Offline-Speicherung", Inches(3.4), RGBColor(219, 234, 254)),
        ("Firebase", "Cloud-Sync für Kommentare", Inches(4.6), RGBColor(220, 252, 231)),
        ("React UI", "Interaktive Dashboards", Inches(5.8), RGBColor(243, 232, 255))
    ]

    for title, desc, y_pos, bg_color in layers:
        # Layer-Box
        layer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.5), y_pos, Inches(5.5), Inches(1))
        layer.fill.solid()
        layer.fill.fore_color.rgb = bg_color
        layer.line.color.rgb = GRAY_LIGHT

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.7), y_pos + Inches(0.1),
                                            Inches(5), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = GRAY_DARK

        # Description
        desc_box = slide.shapes.add_textbox(Inches(0.7), y_pos + Inches(0.5),
                                           Inches(5), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY_LIGHT

    # Vorteile rechts
    benefits_title = slide.shapes.add_textbox(Inches(7), Inches(2.0), Inches(5), Inches(0.5))
    tf = benefits_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Vorteile der Architektur"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE_PRIMARY

    benefits = [
        ("Offline-Fähig", "Daten immer verfügbar, auch ohne Internet"),
        ("Schnelle Performance", "Lokale Daten = blitzschnelle Abfragen"),
        ("Cloud-Backup", "Kommentare in Firebase gesichert"),
        ("Skalierbar", "Große Datenmengen problemlos möglich"),
        ("Sicher", "Keine sensiblen Daten in der Cloud")
    ]

    for i, (title, desc) in enumerate(benefits):
        # Benefit card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(7), Inches(2.6 + i * 0.85),
                                     Inches(5.5), Inches(0.75))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(248, 250, 252)
        card.line.color.rgb = BLUE_SECONDARY

        title_box = slide.shapes.add_textbox(Inches(7.2), Inches(2.65 + i * 0.85),
                                            Inches(5), Inches(0.35))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BLUE_PRIMARY

        desc_box = slide.shapes.add_textbox(Inches(7.2), Inches(2.95 + i * 0.85),
                                           Inches(5), Inches(0.35))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = GRAY_LIGHT


def add_import_slide(prs):
    """Import-System-Folie"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_header(slide, "Excel-Import System")

    # Prozess-Schritte
    steps = [
        ("1", "Upload", "Multi-File Drag & Drop"),
        ("2", "Validierung", "Schema-Prüfung"),
        ("3", "Mapping", "Spalten-Zuordnung"),
        ("4", "Import", "IndexedDB Speicherung"),
        ("5", "Fertig", "Dashboard aktualisiert")
    ]

    for i, (num, title, desc) in enumerate(steps):
        x_pos = 0.5 + i * 2.5

        # Nummer-Kreis
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       Inches(x_pos + 0.6), Inches(2.0),
                                       Inches(0.8), Inches(0.8))
        circle.fill.solid()
        circle.fill.fore_color.rgb = BLUE_PRIMARY
        circle.line.fill.background()

        num_box = slide.shapes.add_textbox(Inches(x_pos + 0.7), Inches(2.1),
                                          Inches(0.6), Inches(0.6))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Titel
        title_box = slide.shapes.add_textbox(Inches(x_pos), Inches(3.0),
                                            Inches(2), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = GRAY_DARK
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc_box = slide.shapes.add_textbox(Inches(x_pos), Inches(3.4),
                                           Inches(2), Inches(0.5))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY_LIGHT
        p.alignment = PP_ALIGN.CENTER

        # Pfeil (außer beim letzten)
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                          Inches(x_pos + 2), Inches(2.2),
                                          Inches(0.4), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = BLUE_SECONDARY
            arrow.line.fill.background()

    # Features unten
    features_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(0.4))
    tf = features_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Import-Features:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GRAY_DARK

    features = [
        ("Automatische Spalten-Erkennung", "System erkennt Datentypen automatisch"),
        ("Fortschrittsanzeige", "Visuelle Rückmeldung während des Imports"),
        ("Validierungsergebnisse", "Klare Fehler- und Warnmeldungen"),
        ("Multi-Department Support", "Sales, Produktion & Projekte in einem Import")
    ]

    for i, (title, desc) in enumerate(features):
        col = i % 2
        row = i // 2

        title_box = slide.shapes.add_textbox(Inches(0.5 + col * 6.3), Inches(5.0 + row * 0.9),
                                            Inches(5.5), Inches(0.35))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "✓ " + title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BLUE_PRIMARY

        desc_box = slide.shapes.add_textbox(Inches(0.7 + col * 6.3), Inches(5.3 + row * 0.9),
                                           Inches(5.3), Inches(0.35))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY_LIGHT


def add_summary_slide(prs):
    """Zusammenfassung-Folie"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_header(slide, "Zusammenfassung")

    # Hauptvorteile
    benefits = [
        ("Integrierte Lösung", "Ein System für Sales, Produktion und Projektmanagement"),
        ("Offline-First", "Zuverlässige Performance, auch ohne Internetverbindung"),
        ("Echtzeit-Sync", "Firebase-basierte Kommentar-Synchronisation im Team"),
        ("Excel-Integration", "Nahtloser Import bestehender Daten"),
        ("Moderne UI", "Intuitive Benutzeroberfläche mit TailwindCSS"),
        ("Erweiterbar", "Modulare Architektur für zukünftige Features")
    ]

    for i, (title, desc) in enumerate(benefits):
        col = i % 2
        row = i // 2

        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.5 + col * 6.3), Inches(1.6 + row * 1.4),
                                     Inches(5.8), Inches(1.2))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(248, 250, 252)
        card.line.color.rgb = BLUE_SECONDARY

        # Icon
        icon = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(0.7 + col * 6.3), Inches(1.8 + row * 1.4),
                                     Inches(0.5), Inches(0.5))
        icon.fill.solid()
        icon.fill.fore_color.rgb = BLUE_PRIMARY
        icon.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(1.4 + col * 6.3), Inches(1.7 + row * 1.4),
                                            Inches(4.5), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLUE_PRIMARY

        # Description
        desc_box = slide.shapes.add_textbox(Inches(1.4 + col * 6.3), Inches(2.1 + row * 1.4),
                                           Inches(4.5), Inches(0.5))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY_DARK

    # Technologie-Summary
    tech_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12), Inches(0.4))
    tf = tech_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Tech-Stack: React + TypeScript + TailwindCSS + IndexedDB + Firebase + Recharts"
    p.font.size = Pt(14)
    p.font.color.rgb = GRAY_LIGHT
    p.alignment = PP_ALIGN.CENTER


def add_thank_you_slide(prs):
    """Danke-Folie"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Hintergrund
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE_PRIMARY
    shape.line.fill.background()

    # Logo-Bereich
    icon_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.9), Inches(1.5), Inches(1.5), Inches(1.5))
    icon_shape.fill.solid()
    icon_shape.fill.fore_color.rgb = WHITE
    icon_shape.line.fill.background()

    # Danke Text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.33), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Vielen Dank!"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Untertitel
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.33), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Fragen & Diskussion"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(191, 219, 254)
    p.alignment = PP_ALIGN.CENTER

    # Projekt-Name
    project_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(12.33), Inches(0.5))
    tf = project_box.text_frame
    p = tf.paragraphs[0]
    p.text = "GALADRIEL - Business Intelligence Dashboard"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(147, 197, 253)
    p.alignment = PP_ALIGN.CENTER


def create_presentation():
    """Hauptfunktion zum Erstellen der Präsentation"""
    prs = Presentation()

    # Widescreen 16:9 Format
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Folien erstellen
    add_title_slide(prs)
    add_agenda_slide(prs)
    add_overview_slide(prs)
    add_tech_stack_slide(prs)
    add_sales_module_slide(prs)
    add_controlling_module_slide(prs)
    add_data_comparison_slide(prs)
    add_architecture_slide(prs)
    add_import_slide(prs)
    add_summary_slide(prs)
    add_thank_you_slide(prs)

    # Speichern
    output_path = "/home/user/Neubau-Claude/Galadriel_Praesentation.pptx"
    prs.save(output_path)
    print(f"✅ Präsentation erfolgreich erstellt: {output_path}")
    print(f"📊 Anzahl der Folien: {len(prs.slides)}")
    return output_path


if __name__ == "__main__":
    create_presentation()
